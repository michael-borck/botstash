"""PPTX text extraction with slide numbers and speaker notes."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def extract_pptx(file_path: Path) -> str:
    """Extract plain text from a PPTX file, one section per slide."""
    prs = Presentation(str(file_path))
    sections: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"--- Slide {i} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker Notes] {notes}")

        sections.append("\n".join(parts))

    return "\n\n".join(sections)
