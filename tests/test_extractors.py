"""Tests for content extractors."""

from pathlib import Path

import docx
from pptx import Presentation

from botstash.extractors import extract_file
from botstash.extractors.docx import extract_docx
from botstash.extractors.pptx import extract_pptx
from botstash.extractors.qti import extract_qti
from botstash.extractors.unit_outline import extract_unit_outline
from botstash.extractors.url_tracker import extract_urls, log_urls
from botstash.extractors.vtt import extract_vtt


def test_extract_vtt(tmp_path: Path) -> None:
    """VTT extraction strips timestamps and returns clean text."""
    vtt_content = (
        "WEBVTT\n\n"
        "00:00:01.000 --> 00:00:03.000\n"
        "Hello world\n\n"
        "00:00:04.000 --> 00:00:06.000\n"
        "This is a test\n"
    )
    vtt_file = tmp_path / "test.vtt"
    vtt_file.write_text(vtt_content)

    result = extract_vtt(vtt_file)
    assert "Hello world" in result
    assert "This is a test" in result
    assert "-->" not in result
    assert "00:00" not in result


def test_extract_docx(tmp_path: Path) -> None:
    """DOCX extraction preserves headings as markdown markers."""
    doc = docx.Document()
    doc.add_heading("Main Title", level=1)
    doc.add_paragraph("Some body text here.")
    doc.add_heading("Sub Section", level=2)
    doc.add_paragraph("More text.")
    doc_path = tmp_path / "test.docx"
    doc.save(str(doc_path))

    result = extract_docx(doc_path)
    assert "# Main Title" in result
    assert "## Sub Section" in result
    assert "Some body text here." in result
    assert "More text." in result


def test_extract_pptx(tmp_path: Path) -> None:
    """PPTX extraction includes slide numbers and speaker notes."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Slide One Title"
    slide1.placeholders[1].text = "Slide one content"

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide Two Title"
    slide2.placeholders[1].text = "Slide two content"
    slide2.notes_slide.notes_text_frame.text = "These are speaker notes"

    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))

    result = extract_pptx(pptx_path)
    assert "--- Slide 1 ---" in result
    assert "--- Slide 2 ---" in result
    assert "Slide One Title" in result
    assert "Slide two content" in result
    assert "[Speaker Notes] These are speaker notes" in result


def test_extract_qti_v12(tmp_path: Path) -> None:
    """QTI 1.2 extraction returns question text only."""
    qti_xml = """<?xml version="1.0" encoding="UTF-8"?>
    <questestinterop>
      <assessment>
        <section>
          <item>
            <presentation>
              <mattext>What is the capital of France?</mattext>
            </presentation>
            <resprocessing>
              <respcondition>
                <conditionvar><varequal>Paris</varequal></conditionvar>
              </respcondition>
            </resprocessing>
          </item>
          <item>
            <presentation>
              <mattext>What is 2 + 2?</mattext>
            </presentation>
          </item>
        </section>
      </assessment>
    </questestinterop>"""
    qti_path = tmp_path / "quiz.xml"
    qti_path.write_text(qti_xml)

    result = extract_qti(qti_path)
    assert "What is the capital of France?" in result
    assert "What is 2 + 2?" in result
    assert "Paris" not in result


def test_extract_urls() -> None:
    """URL extraction finds href and src URLs from HTML."""
    html = """
    <a href="https://example.com/video">Watch</a>
    <iframe src="https://echo360.org/lesson/123"></iframe>
    <a href="/relative/path">Internal</a>
    <img src="https://cdn.example.com/image.png" />
    """
    urls = extract_urls(html)
    assert "https://example.com/video" in urls
    assert "https://echo360.org/lesson/123" in urls
    assert "https://cdn.example.com/image.png" in urls
    assert "/relative/path" not in urls  # relative URLs excluded


def test_log_urls(tmp_path: Path) -> None:
    """URL logger appends to file with context."""
    log_path = tmp_path / "urls_log.txt"
    log_urls(
        ["https://example.com/video1"],
        "Week 1 Lecture",
        log_path,
    )
    log_urls(
        ["https://example.com/video2"],
        "Week 2 Lecture",
        log_path,
    )
    content = log_path.read_text()
    assert "[Week 1 Lecture] https://example.com/video1" in content
    assert "[Week 2 Lecture] https://example.com/video2" in content


def test_log_urls_empty(tmp_path: Path) -> None:
    """URL logger does nothing with an empty URL list."""
    log_path = tmp_path / "urls_log.txt"
    log_urls([], "context", log_path)
    assert not log_path.exists()


def test_unit_outline_docx(tmp_path: Path) -> None:
    """Unit outline delegates to DOCX extractor."""
    doc = docx.Document()
    doc.add_heading("Unit Outline", level=1)
    doc.add_paragraph("Learning outcomes for this unit.")
    doc_path = tmp_path / "outline.docx"
    doc.save(str(doc_path))

    result = extract_unit_outline(str(doc_path))
    assert "Unit Outline" in result
    assert "Learning outcomes" in result


def test_extract_file_dispatcher(tmp_path: Path) -> None:
    """Registry dispatcher routes by file extension."""
    vtt_content = (
        "WEBVTT\n\n"
        "00:00:01.000 --> 00:00:03.000\n"
        "Test caption\n"
    )
    vtt_file = tmp_path / "test.vtt"
    vtt_file.write_text(vtt_content)

    result = extract_file(vtt_file)
    assert result is not None
    assert "Test caption" in result


def test_extract_file_unknown_extension(tmp_path: Path) -> None:
    """Registry returns None for unknown file types."""
    txt_file = tmp_path / "readme.txt"
    txt_file.write_text("hello")

    result = extract_file(txt_file)
    assert result is None
